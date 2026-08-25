"""Tests for figure extraction and the image tools."""

import io
from unittest.mock import AsyncMock, patch

import pytest
from fastmcp import FastMCP

import canvas_mcp.core.config as config_module
from canvas_mcp.core.images import extract_document_images, save_figures
from canvas_mcp.tools.images import register_image_tools


def _png(size=(64, 64), color=(200, 30, 30)):
    from PIL import Image
    buf = io.BytesIO()
    Image.new("RGB", size, color).save(buf, format="PNG")
    return buf.getvalue()


def _noisy_png(size=(128, 128)):
    """Incompressible image that clears the tiny-logo byte filter."""
    import os

    from PIL import Image
    img = Image.frombytes("RGB", size, os.urandom(size[0] * size[1] * 3))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _pptx_with_pictures(*pngs):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    for png in pngs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(io.BytesIO(png), Inches(1), Inches(1))
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def fake_root(tmp_path, monkeypatch):
    root = tmp_path / "hulms-agent"
    root.mkdir()
    monkeypatch.setattr(config_module, "REPO_ROOT", root)
    (tmp_path / "spaces" / "c1").mkdir(parents=True)
    yield tmp_path


def test_pptx_pictures_extracted():
    data = _pptx_with_pictures(_png())
    result = extract_document_images(data, "deck.pptx", min_bytes=0)
    assert len(result["images"]) == 1
    name, blob = result["images"][0]
    assert name.startswith("slide001") and blob[:4] == b"\x89PNG"
    assert result["unit"] == "slide" and result["total"] == 1


def test_duplicate_pictures_deduped():
    same = _png()
    data = _pptx_with_pictures(same, same)
    result = extract_document_images(data, "deck.pptx", min_bytes=0)
    assert len(result["images"]) == 1


def test_tiny_images_filtered():
    data = _pptx_with_pictures(_png(size=(4, 4)))
    result = extract_document_images(data, "deck.pptx", min_bytes=4096)
    assert result["images"] == [] and result["capped"] is False


def test_unsupported_format_returns_note():
    result = extract_document_images(b"x", "notes.docx")
    assert "No image extractor" in result["error"]


def test_standalone_image_passthrough():
    png = _png()
    result = extract_document_images(png, "diagram.png", min_bytes=0)
    assert result["images"] == [("image.png", png)]


def test_slide_window_targets_later_slides():
    """The textbook problem: early art must not shadow a targeted window."""
    early, wanted = _noisy_png(), _noisy_png()
    data = _pptx_with_pictures(early, early, wanted)
    result = extract_document_images(data, "deck.pptx", min_bytes=0, first=3, last=3)
    assert len(result["images"]) == 1
    assert result["images"][0][0].startswith("slide003")
    assert result["scannedFrom"] == 3 and result["scannedTo"] == 3


def test_cap_is_reported_not_silent():
    data = _pptx_with_pictures(_noisy_png(), _noisy_png(), _noisy_png())
    result = extract_document_images(data, "deck.pptx", min_bytes=0, max_images=1)
    assert len(result["images"]) == 1
    assert result["capped"] is True
    assert result["scannedTo"] < result["total"]


def test_save_figures_returns_view_and_embed(fake_root):
    saved = save_figures("canvas-42", [("page001-1.png", _png())])
    assert saved[0]["file"].endswith("page001-1.png")
    assert saved[0]["embed"].startswith("/api/spacefile?p=.figures/canvas-42/")
    from pathlib import Path
    assert Path(saved[0]["file"]).exists()


# ------------------------------------------------------------------- tools

def get_tool(name):
    mcp = FastMCP("test")
    captured = {}
    original = mcp.tool

    def capturing(*a, **k):
        d = original(*a, **k)
        def w(fn):
            captured[fn.__name__] = fn
            return d(fn)
        return w
    mcp.tool = capturing
    register_image_tools(mcp)
    return captured[name]


async def test_document_images_requires_exactly_one_source():
    tool = get_tool("get_document_images")
    assert "error" in await tool()
    assert "error" in await tool(file_id=1, local_name="x.pdf")


async def test_document_images_from_local_pptx(fake_root):
    deck = fake_root / "spaces" / "c1" / "deck.pptx"
    deck.write_bytes(_pptx_with_pictures(_noisy_png()))
    result = await get_tool("get_document_images")(local_name="deck.pptx")
    assert result["count"] == 1
    assert result["images"][0]["embed"].startswith("/api/spacefile?p=.figures/")


async def test_document_images_canvas_error_propagates():
    with patch("canvas_mcp.tools.images.fetch_file_bytes",
               new=AsyncMock(return_value={"error": "Could not read file 9"})):
        assert "error" in await get_tool("get_document_images")(file_id=9)


async def test_fetch_web_image_saves_png(fake_root):
    png = _png()

    class FakeResp:
        headers = {"content-type": "image/png"}
        content = png
        def raise_for_status(self): pass

    class FakeClient:
        def __init__(self, **kw): pass
        async def __aenter__(self): return self
        async def __aexit__(self, *a): return False
        async def get(self, url): return FakeResp()

    with patch("canvas_mcp.tools.images.httpx.AsyncClient", FakeClient):
        result = await get_tool("fetch_web_image")("https://example.com/fig.png")
    assert result["embed"].startswith("/api/spacefile?p=.figures/web/")
    assert result["sourceUrl"] == "https://example.com/fig.png"


async def test_fetch_web_image_rejects_non_image(fake_root):
    class FakeResp:
        headers = {"content-type": "text/html"}
        content = b"<html>"
        def raise_for_status(self): pass

    class FakeClient:
        def __init__(self, **kw): pass
        async def __aenter__(self): return self
        async def __aexit__(self, *a): return False
        async def get(self, url): return FakeResp()

    with patch("canvas_mcp.tools.images.httpx.AsyncClient", FakeClient):
        assert "error" in await get_tool("fetch_web_image")("https://example.com/page")


def test_render_pdf_pages_rasterizes_vector_content():
    """Vector-only PDFs (textbook line-art) yield nothing to the embedded-
    image extractor but render fine as whole pages."""
    from canvas_mcp.core.images import render_pdf_pages
    from tests.core.test_extract import MINI_PDF  # text-only, zero embedded images

    embedded = extract_document_images(MINI_PDF, "book.pdf", min_bytes=0)
    assert embedded["images"] == []

    rendered = render_pdf_pages(MINI_PDF, 1, 1)
    assert len(rendered["images"]) == 1
    name, blob = rendered["images"][0]
    assert name == "render-page001.png" and blob[:4] == b"\x89PNG"


def test_render_window_larger_than_document_is_not_capped():
    from canvas_mcp.core.images import render_pdf_pages
    from tests.core.test_extract import MINI_PDF
    result = render_pdf_pages(MINI_PDF, 1, 50)
    assert result["scannedTo"] == 1  # document only has one page
    assert result["capped"] is False  # nothing was cut off — honest


def test_render_caps_pages_per_call():
    from pypdf import PdfWriter

    from canvas_mcp.core.images import MAX_RENDER_PAGES, render_pdf_pages
    writer = PdfWriter()
    for _ in range(10):
        writer.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    writer.write(buf)
    result = render_pdf_pages(buf.getvalue(), 1, 10)
    assert len(result["images"]) == MAX_RENDER_PAGES
    assert result["scannedTo"] == MAX_RENDER_PAGES
    assert result["capped"] is True


def test_render_beyond_document_errors():
    from canvas_mcp.core.images import render_pdf_pages
    from tests.core.test_extract import MINI_PDF
    assert "error" in render_pdf_pages(MINI_PDF, 99, 100)


def test_crop_image_file_cuts_the_right_region(tmp_path):
    from PIL import Image

    from canvas_mcp.core.images import crop_image_file
    img = Image.new("RGB", (200, 100), (0, 0, 255))
    for x in range(100, 200):
        for y in range(50, 100):
            img.putpixel((x, y), (255, 0, 0))  # bottom-right quadrant red
    src = tmp_path / "page.png"
    img.save(src)

    name, blob = crop_image_file(src, 50, 50, 100, 100)
    out = Image.open(io.BytesIO(blob))
    assert out.size == (100, 50)
    assert out.getpixel((10, 10)) == (255, 0, 0)
    assert "crop-50x50-100x100" in name


def test_crop_image_file_rejects_bad_box(tmp_path):
    from PIL import Image

    from canvas_mcp.core.images import crop_image_file
    src = tmp_path / "x.png"
    Image.new("RGB", (10, 10)).save(src)
    assert isinstance(crop_image_file(src, 60, 0, 40, 100), str)  # left >= right
    assert isinstance(crop_image_file(src, 0, 0, 100, 101), str)  # out of range


async def test_crop_tool_containment(fake_root):
    result = await get_tool("crop_image")("C:/Windows/win.ini", 0, 0, 50, 50)
    assert "error" in result


async def test_crop_tool_roundtrip(fake_root):
    from canvas_mcp.core.images import save_figures
    saved = save_figures("canvas-7-pages", [("render-page001.png", _noisy_png((100, 100)))])
    result = await get_tool("crop_image")(saved[0]["file"], 0, 0, 50, 50)
    assert result["embed"].startswith("/api/spacefile?p=.figures/canvas-7-pages/")
    from pathlib import Path
    assert Path(result["file"]).exists()
