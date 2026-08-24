"""Extraction pipeline tests: real generated files, no mocks."""

import io

from canvas_mcp.core.extract import extract_text, is_extractable

# A minimal but complete one-page PDF with one text object.
MINI_PDF = b"""%PDF-1.4
1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj
2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj
3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >> endobj
4 0 obj << /Length 60 >> stream
BT /F1 12 Tf 72 712 Td (Recursion uses a call stack) Tj ET
endstream endobj
5 0 obj << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> endobj
xref
0 1
0000000000 65535 f
trailer << /Root 1 0 R /Size 6 >>
startxref
0
%%EOF
"""


def test_pdf_text_extraction():
    result = extract_text(MINI_PDF, "notes.pdf")
    assert result.status == "ok"
    assert "Recursion uses a call stack" in result.text


def test_scanned_pdf_detected_and_skipped():
    from pypdf import PdfWriter

    writer = PdfWriter()
    for _ in range(3):
        writer.add_blank_page(width=612, height=792)
    buf = io.BytesIO()
    writer.write(buf)

    result = extract_text(buf.getvalue(), "scan.pdf")
    assert result.status == "scanned"
    assert result.text == ""
    assert "Skipped" in result.note


def test_pptx_extraction_includes_slides_and_notes():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Queues"
    slide.placeholders[1].text = "FIFO: enqueue at rear, dequeue at front"
    slide.notes_slide.notes_text_frame.text = "mention circular queues"
    buf = io.BytesIO()
    prs.save(buf)

    result = extract_text(buf.getvalue(), "week3.pptx")
    assert result.status == "ok"
    assert "FIFO" in result.text
    assert "circular queues" in result.text
    assert "slide 1" in result.text


def test_docx_extraction_includes_tables():
    from docx import Document

    doc = Document()
    doc.add_paragraph("Stack operations")
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "push"
    table.rows[0].cells[1].text = "O(1)"
    buf = io.BytesIO()
    doc.save(buf)

    result = extract_text(buf.getvalue(), "handout.docx")
    assert result.status == "ok"
    assert "Stack operations" in result.text
    assert "push" in result.text and "O(1)" in result.text


def test_html_is_stripped():
    result = extract_text(b"<h1>Heaps</h1><p>parent &lt; children</p>", "page.html")
    assert result.status == "ok"
    assert "Heaps" in result.text and "parent < children" in result.text
    assert "<p>" not in result.text


def test_notebook_cells_extracted():
    nb = b'{"cells": [{"cell_type": "markdown", "source": ["# Sorting"]}, {"cell_type": "code", "source": ["def bubble(a):\\n", "    pass"]}]}'
    result = extract_text(nb, "lab.ipynb")
    assert result.status == "ok"
    assert "# Sorting" in result.text and "def bubble" in result.text


def test_legacy_office_reported_unsupported():
    result = extract_text(b"\xd0\xcf\x11\xe0junk", "old.ppt")
    assert result.status == "unsupported"
    assert result.note


def test_video_is_not_extractable():
    assert not is_extractable("lecture.mov")
    assert is_extractable("slides.pdf")


def test_extensionless_display_name_falls_back_to_mime():
    # Measured live: Canvas display_name often drops the extension.
    assert is_extractable("Designing with Stacks", None, "application/pdf")
    assert is_extractable("Notes", "notes.docx", None)
    assert not is_extractable("Mystery blob", None, "application/octet-stream")
    result = extract_text(MINI_PDF, "Designing with Stacks", None, "application/pdf")
    assert result.status == "ok"
    assert "call stack" in result.text


def test_corrupt_file_reports_error_not_crash():
    result = extract_text(b"this is not a pdf", "broken.pdf")
    assert result.status in ("error", "scanned")
