"""Figure extraction: pull embedded images out of course documents.

Lecture slides carry their meaning in diagrams as much as text. Extracted
figures are saved under spaces/.figures/ where the coach can view them with
its (vision-capable) Read tool and the UI can serve them inline.
"""

import hashlib
import io
from pathlib import Path
from typing import Any

from .extract import file_extension
from .local_files import spaces_root

# Below this, an image is almost certainly a logo, bullet glyph, or border.
MIN_IMAGE_BYTES = 4096
DEFAULT_MAX_IMAGES = 12

_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"}


def _dedupe_key(data: bytes) -> str:
    return hashlib.md5(data).hexdigest()


def _pdf_images(
    data: bytes, cap: int, min_bytes: int,
    first: int = 1, last: int | None = None,
) -> dict[str, Any]:
    from pypdf import PdfReader

    out: list[tuple[str, bytes]] = []
    seen: set[str] = set()
    reader = PdfReader(io.BytesIO(data))
    total = len(reader.pages)
    first = max(1, first)
    last = min(last or total, total)
    reached = first
    capped = False
    for page_num in range(first, last + 1):
        reached = page_num
        try:
            images = reader.pages[page_num - 1].images
        except Exception:
            continue  # one bad page must not kill the rest
        for img in images:
            try:
                blob = img.data
            except Exception:
                continue
            if len(blob) < min_bytes:
                continue
            key = _dedupe_key(blob)
            if key in seen:
                continue
            seen.add(key)
            ext = file_extension(img.name) or ".png"
            out.append((f"page{page_num:03d}-{len(out) + 1}{ext}", blob))
            if len(out) >= cap:
                capped = page_num < last
                return {"images": out, "unit": "page", "total": total,
                        "scannedFrom": first, "scannedTo": page_num, "capped": capped}
    return {"images": out, "unit": "page", "total": total,
            "scannedFrom": first, "scannedTo": reached, "capped": False}


def _pptx_images(
    data: bytes, cap: int, min_bytes: int,
    first: int = 1, last: int | None = None,
) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    out: list[tuple[str, bytes]] = []
    seen: set[str] = set()
    prs = Presentation(io.BytesIO(data))
    slides = list(prs.slides)
    total = len(slides)
    first = max(1, first)
    last = min(last or total, total)
    reached = first
    for slide_num in range(first, last + 1):
        reached = slide_num
        for shape in slides[slide_num - 1].shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                blob = shape.image.blob
                ext = "." + shape.image.ext
            except Exception:
                continue
            if len(blob) < min_bytes:
                continue
            key = _dedupe_key(blob)
            if key in seen:
                continue
            seen.add(key)
            out.append((f"slide{slide_num:03d}-{len(out) + 1}{ext}", blob))
            if len(out) >= cap:
                return {"images": out, "unit": "slide", "total": total,
                        "scannedFrom": first, "scannedTo": slide_num,
                        "capped": slide_num < last}
    return {"images": out, "unit": "slide", "total": total,
            "scannedFrom": first, "scannedTo": reached, "capped": False}


def extract_document_images(
    data: bytes,
    filename: str,
    max_images: int = DEFAULT_MAX_IMAGES,
    min_bytes: int = MIN_IMAGE_BYTES,
    first: int = 1,
    last: int | None = None,
) -> dict[str, Any]:
    """Embedded images of a PDF/PPTX, optionally within a page/slide window.

    Returns {images: [(name, bytes)], unit, total, scannedFrom, scannedTo,
    capped} — or {error: str}. `capped` means the max_images limit stopped
    the scan before the window's end: later figures exist but were not
    reached (the textbook problem — front matter exhausts the cap).
    """
    ext = file_extension(filename)
    try:
        if ext == ".pdf":
            return _pdf_images(data, max_images, min_bytes, first, last)
        if ext == ".pptx":
            return _pptx_images(data, max_images, min_bytes, first, last)
        if ext in _IMAGE_EXTS:
            # The document IS an image.
            images = [(f"image{ext}", data)] if len(data) >= min_bytes else []
            return {"images": images, "unit": "image", "total": 1,
                    "scannedFrom": 1, "scannedTo": 1, "capped": False}
        return {"error": f"No image extractor for '{ext or 'file without extension'}' (PDF and PPTX carry figures)."}
    except Exception as e:
        return {"error": f"Image extraction failed: {type(e).__name__}: {e}"}


MAX_RENDER_PAGES = 6


def render_pdf_pages(
    data: bytes, first: int, last: int, scale: float = 2.0
) -> dict[str, Any]:
    """Rasterize whole PDF pages to PNG — the answer for VECTOR figures.

    Textbook diagrams are usually line-art, not embedded images, so
    extract_document_images finds nothing; rendering the page itself makes
    them visible regardless. Capped at MAX_RENDER_PAGES per call.
    """
    try:
        import pypdfium2 as pdfium
    except ImportError:
        return {"error": "pypdfium2 is not installed."}

    try:
        pdf = pdfium.PdfDocument(data)
    except Exception as e:
        return {"error": f"Could not open PDF: {type(e).__name__}: {e}"}
    try:
        total = len(pdf)
        first = max(1, first)
        last = min(last, total)
        if first > total:
            return {"error": f"Page {first} is beyond the document ({total} pages)."}
        capped = (last - first + 1) > MAX_RENDER_PAGES
        if capped:
            last = first + MAX_RENDER_PAGES - 1

        out: list[tuple[str, bytes]] = []
        for page_num in range(first, last + 1):
            page = pdf[page_num - 1]
            bitmap = page.render(scale=scale)
            buf = io.BytesIO()
            bitmap.to_pil().save(buf, format="PNG")
            out.append((f"render-page{page_num:03d}.png", buf.getvalue()))
        return {"images": out, "unit": "page", "total": total,
                "scannedFrom": first, "scannedTo": last, "capped": capped}
    finally:
        pdf.close()


def crop_image_file(
    src: Path, left: float, top: float, right: float, bottom: float
) -> tuple[str, bytes] | str:
    """Crop an image by percentage box (0-100 from top-left); str = error.

    Percentages, not pixels, because the caller is a model looking at the
    image — "the figure spans roughly 10-90% wide, 20-55% down" is how
    vision naturally localizes.
    """
    from PIL import Image

    if not (0 <= left < right <= 100 and 0 <= top < bottom <= 100):
        return "Box must satisfy 0 <= left < right <= 100 and 0 <= top < bottom <= 100."
    try:
        img = Image.open(src)
    except Exception as e:
        return f"Could not open image: {type(e).__name__}: {e}"
    w, h = img.size
    box = (
        int(w * left / 100), int(h * top / 100),
        int(w * right / 100), int(h * bottom / 100),
    )
    cropped = img.crop(box)
    buf = io.BytesIO()
    cropped.save(buf, format="PNG")
    name = f"{src.stem}-crop-{int(left)}x{int(top)}-{int(right)}x{int(bottom)}.png"
    return name, buf.getvalue()


def figures_dir(source_key: str) -> Path:
    """Where a source's extracted figures live (under the gitignored spaces/)."""
    safe = "".join(c if c.isalnum() or c in "-_." else "_" for c in source_key)[:80]
    return spaces_root() / ".figures" / safe


def save_figures(source_key: str, images: list[tuple[str, bytes]]) -> list[dict[str, Any]]:
    """Write figures to disk; returns view/embed handles for each.

    `file` is the absolute path (for the coach's vision Read); `embed` is the
    UI-served URL path for showing the image inline in chat.
    """
    directory = figures_dir(source_key)
    directory.mkdir(parents=True, exist_ok=True)
    root = spaces_root()
    out = []
    for name, blob in images:
        target = directory / name
        target.write_bytes(blob)
        rel = target.relative_to(root).as_posix()
        out.append({
            "name": name,
            "file": str(target),
            "embed": f"/api/spacefile?p={rel}",
            "sizeBytes": len(blob),
        })
    return out
