"""File text extraction: PDF, PPTX, DOCX, and plain-text-ish formats.

Everything runs in-process (pypdf, python-pptx, python-docx). Scanned PDFs
are detected and skipped rather than half-OCRed — a page-image PDF yields
almost no text, and returning that sliver as if it were the document would
be a confident lie. Status is always reported so callers can say honestly
what was and wasn't read.
"""

import io
import json
import os
from dataclasses import dataclass

from .text import strip_html_tags

# Formats worth downloading at all. Anything else is skipped BEFORE download.
PDF_EXTS = {".pdf"}
PPTX_EXTS = {".pptx"}
DOCX_EXTS = {".docx"}
TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".py", ".java", ".c", ".cpp", ".h"}
HTML_EXTS = {".html", ".htm"}
NOTEBOOK_EXTS = {".ipynb"}

EXTRACTABLE_EXTS = PDF_EXTS | PPTX_EXTS | DOCX_EXTS | TEXT_EXTS | HTML_EXTS | NOTEBOOK_EXTS

# Legacy Office formats (.ppt/.doc) need external converters — reported as
# unsupported rather than silently absent.
UNSUPPORTED_KNOWN = {".ppt", ".doc", ".xls", ".xlsx"}

# A "PDF with pages but almost no text" is a scan. Threshold is deliberately
# low: real slide decks average hundreds of chars per page.
SCANNED_AVG_CHARS_PER_PAGE = 25

MAX_DOWNLOAD_BYTES = 80 * 1024 * 1024


@dataclass
class Extraction:
    status: str  # "ok" | "scanned" | "unsupported" | "error"
    text: str = ""
    note: str = ""


def file_extension(name: str | None) -> str:
    return os.path.splitext(name or "")[1].lower()


# Canvas display_name frequently drops the extension ("Designing with Stacks"
# for a PDF — measured live, 27 such files in one course), so extension
# routing needs the MIME type as a fallback.
_CONTENT_TYPE_EXTS = {
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "text/plain": ".txt",
    "text/markdown": ".md",
    "text/html": ".html",
    "text/csv": ".csv",
    "application/json": ".json",
    "text/x-python": ".py",
}


def effective_extension(
    display_name: str | None,
    filename: str | None = None,
    content_type: str | None = None,
) -> str:
    """Best-effort extension: display_name, then filename, then MIME type."""
    for name in (display_name, filename):
        ext = file_extension(name)
        if ext:
            return ext
    return _CONTENT_TYPE_EXTS.get((content_type or "").split(";")[0].strip(), "")


def is_extractable(
    name: str | None,
    filename: str | None = None,
    content_type: str | None = None,
) -> bool:
    return effective_extension(name, filename, content_type) in EXTRACTABLE_EXTS


def _extract_pdf(data: bytes) -> Extraction:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = [page.extract_text() or "" for page in reader.pages]
    text = "\n\n".join(p.strip() for p in pages if p.strip())
    if pages and len(text) < SCANNED_AVG_CHARS_PER_PAGE * len(pages):
        return Extraction(
            "scanned",
            note=f"PDF has {len(pages)} pages but almost no extractable text; "
            "likely a scan. Skipped rather than half-OCRed.",
        )
    return Extraction("ok", text)


def _extract_pptx(data: bytes) -> Extraction:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    chunks: list[str] = []
    for idx, slide in enumerate(prs.slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    parts.append(t)
        if getattr(slide, "has_notes_slide", False) and slide.notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[speaker notes] {notes}")
        if parts:
            chunks.append(f"--- slide {idx} ---\n" + "\n".join(parts))
    return Extraction("ok", "\n\n".join(chunks))


def _extract_docx(data: bytes) -> Extraction:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append("\t".join(cells))
    return Extraction("ok", "\n".join(parts))


def _extract_notebook(data: bytes) -> Extraction:
    nb = json.loads(data.decode("utf-8", errors="replace"))
    chunks = []
    for cell in nb.get("cells", []):
        src = "".join(cell.get("source") or [])
        if not src.strip():
            continue
        kind = cell.get("cell_type", "code")
        chunks.append(f"--- {kind} cell ---\n{src.strip()}")
    return Extraction("ok", "\n\n".join(chunks))


def extract_text(
    data: bytes,
    filename: str | None,
    real_filename: str | None = None,
    content_type: str | None = None,
) -> Extraction:
    """Extract plain text from file bytes, routed by effective extension."""
    ext = effective_extension(filename, real_filename, content_type)
    try:
        if ext in PDF_EXTS:
            return _extract_pdf(data)
        if ext in PPTX_EXTS:
            return _extract_pptx(data)
        if ext in DOCX_EXTS:
            return _extract_docx(data)
        if ext in NOTEBOOK_EXTS:
            return _extract_notebook(data)
        if ext in HTML_EXTS:
            return Extraction("ok", strip_html_tags(data.decode("utf-8", errors="replace")))
        if ext in TEXT_EXTS:
            return Extraction("ok", data.decode("utf-8", errors="replace"))
        if ext in UNSUPPORTED_KNOWN:
            return Extraction(
                "unsupported",
                note=f"'{ext}' (legacy Office format) needs an external converter; not extracted.",
            )
        return Extraction("unsupported", note=f"No extractor for '{ext or 'file without extension'}'.")
    except Exception as e:  # extraction must never take the tool down
        return Extraction("error", note=f"Extraction failed: {type(e).__name__}: {e}")
